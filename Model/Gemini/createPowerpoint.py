import os
from pptx import Presentation
from pptx.util import Inches

def createPPT(items):
    prs = Presentation()
    images_folder = os.path.join('Model/Gemini', 'images')
    
    # Get Images
    image_files = [f for f in os.listdir(images_folder) if f.endswith(('.jpg', '.png', '.jpeg'))]

    # Loops items
    for index, item in enumerate(items):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2)).text_frame.text = item
        
        if index < len(image_files):
            slide.shapes.add_picture(os.path.join(images_folder, image_files[index]), Inches(1), Inches(2), width=Inches(8))

    pptx_path = os.path.join('Model/Gemini', 'presentation.pptx')
    prs.save(pptx_path)
    print(f'Presentation created: {pptx_path}')
